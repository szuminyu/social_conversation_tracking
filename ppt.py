from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from datetime import date

from infegy_data import Infegy
from news import Newswhip_article, Newswhip_stats


class COVID_PPT:
    #infegy data
    infegy = Infegy()
    volume_df = infegy.volume()
    state_df = infegy.state()
    dma_df = infegy.dma()
    entity_table = infegy.entity_table()
    #newswhip data
    news = Newswhip_stats()
    trend_df = news.trend_df()
    publisher_df = news.top_publisher()
    news_articles = Newswhip_article()
    article_df = news_articles.engagement_articles()

    # border
    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _set_cell_border(self, cell, border_color="000000", border_width='12700'):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
            ln = self.SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
            solidFill = self.SubElement(ln, 'a:solidFill')
            srgbClr = self.SubElement(solidFill, 'a:srgbClr', val=border_color)
            prstDash = self.SubElement(ln, 'a:prstDash', val='solid')
            round_ = self.SubElement(ln, 'a:round')
            headEnd = self.SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
            tailEnd = self.SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

    #format
    def human_format(self, num):
        magnitude = 0
        while abs(num) >= 1000:
            magnitude += 1
            num /= 1000.0
        # add more suffixes if you need them
        return '%.2f%s' % (num, ['', 'K', 'M', 'G', 'T', 'P'][magnitude])

    #create power point
    def PPT(self):

        prs = Presentation('normal.pptx')

        # P.1 Title Slides
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = 'Coronavirus Media and Social Conversation Tracking'
        subtitle = slide.placeholders[1]
        today = date.today()

        days = ['Monday','Tuesday','Wednesday','Thursday', 'Friday', 'Saturday','Sunday']
        months = ['January', 'February', 'March', 'April', 'May','June','July','August','September', 'October', 'November','December']

        subtitle.text = str(days[date.weekday(today)]) + ', ' + str(months[today.month-1]) + ' ' + str(today.day)+ ', ' + str(today.year)

        txBox1 = slide.shapes.add_textbox(left=Inches(2.23), top=Inches(4.62), width=Inches(9.2), height=Inches(2.22))
        tf = txBox1.text_frame
        p = tf.paragraphs[0]
        p.text = 'This report provides an overview of the social and digital news media conversation surrounding the coronavirus topic. Data in this report comes from: \n' \
                 '(1) NewsWhip, which provides data on articles from over 430,000 website and site feeds, including news sites, blogs, and brand sites \n' \
                 '(2) public social media posts from Facebook, Instagram, Pinterest, and Twitter. \n' \
                 'The following query was used: coronavirus OR “corona virus” OR COVID-19 OR “COVID 19” OR COVID19.'
        p.font.size = Pt(18)
        p.font.bold = False

        ## P.2 Volume line chart ##

        slide_2 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_2.placeholders[0].text = 'Total public social media posts referencing the coronavirus topic (Jan 11 – present)'
        # exclude current day
        self.volume_df = self.volume_df[:-1]
        # make charts
        chart_data = ChartData()
        chart_data.categories = self.volume_df['date'].tolist()
        chart_data.add_series('Volume', self.volume_df['posts_normalized'].tolist())
        x, y, cx, cy = Inches(2), Inches(1.85), Inches(8.9), Inches(4.5)
        chart = slide_2.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart

        #set legend & title
        chart.has_legend = False
        chart.has_title = False
        #set categorical axis (x axis here)
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(12)
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.tick_labels.font.name = "Calibri Light (Headings)"
        category_axis.has_major_gridlines = False
        category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        #set value axis (y axis here)
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(12)
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.tick_labels.font.name = "Calibri Light (Headings)"
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.number_format = '[>999999] #,,"M";#,"K"'

        ## P.3 Article Trend Chart ##
        slide_3 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_3.placeholders[0].text = 'Number of articles referencing the coronavirus topic (Jan 11 – present)'
        # exclude current day
        self.trend_df = self.trend_df[:-1]
        # make chart
        chart_data = ChartData()
        chart_data.categories = self.trend_df['time'].tolist()
        chart_data.add_series('Global', self.trend_df['total_posts_global'].tolist())
        chart_data.add_series('USA', self.trend_df['total_posts_us'].tolist())
        x, y, cx, cy = Inches(1.78), Inches(2), Inches(9.9), Inches(4.5)
        chart = slide_3.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart

        # set legend & title
        chart.has_legend = True
        chart.has_title = False
        # set categorical axis (x axis here)
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(12)
        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.tick_labels.font.name = "Calibri Light (Headings)"
        category_axis.has_major_gridlines = False
        category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        # set value axis (y axis here)
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(12)
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.tick_labels.font.name = "Calibri Light (Headings)"
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.number_format = '[>999999] #,,"M";#,"K"'

        ## P.4 Top 10 US States ##

        slide_4 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_4.placeholders[0].text = 'Top 10 U.S. States by density of social media posts referencing the coronavirus topic (previous 48 hours)'

        ## map
        img_path = 'states.png'
        states_map = slide_4.shapes.add_picture(img_path, left=Inches(0.69), top=Inches(1.76), height=Inches(4.72),
                                                width=Inches(7.08))
        ## table
        table = slide_4.shapes.add_table(11, 2, width=Inches(5), height=Inches(1.88), left=Inches(7.9),
                                                top=Inches(2.31))
        #first row of table
        table.table.cell(0, 0).text = 'Ranking'
        cell = table.table.rows[0].cells[0]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        table.table.cell(0, 1).text = "State"
        cell = table.table.rows[0].cells[1]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        for i in range(1, 11):
            for j in range(0, 2):
                cell = table.table.rows[i].cells[j]
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text = str(self.state_df.iloc[i-1, j])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph = cell.text_frame.paragraphs[0]

        '''
        # P.4 DMA Map and Table
        slide_3 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_3.placeholders[0].text = 'Top 10 U.S. media markets by number of social media posts referencing COVID-19'

        ## map
        img_path = 'dma.png'
        dma_map = slide_3.shapes.add_picture(img_path, left=Inches(0.69), top=Inches(1.76), height=Inches(4.81),
                                                width=Inches(7.26))
        ## table
        dma_table = slide_3.shapes.add_table(11, 3, width=Inches(5), height=Inches(1.88), left=Inches(7.9),
                                                top=Inches(2.31))

        dma_table.table.cell(0, 0).text = 'Ranking'
        cell = dma_table.table.rows[0].cells[0]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        dma_table.table.cell(0, 1).text = "Media Market"
        cell = dma_table.table.rows[0].cells[1]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        dma_table.table.cell(0, 2).text = "Total Posts"
        cell = dma_table.table.rows[0].cells[2]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]


        for i in range(1, 11):
            for j in range(0, 3):
                cell = dma_table.table.rows[i].cells[j]
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text = str(self.dma_df.iloc[i-1,j])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph = cell.text_frame.paragraphs[0]

        ## text box
        txBox1 = slide_3.shapes.add_textbox(left = Inches(0.1), top = Inches(6.86), width = Inches(6), height = Inches(0.86))
        tf = txBox1.text_frame
        p = tf.paragraphs[0]
        p.text = 'Query: coronavirus OR “corona virus” OR COVID-19 OR “COVID 19” OR COVID19 \nPeriod: previous 48 hours'
        p.font.size = Pt(12)
        p.font.bold = False

        '''

        # P.5 Entity Wordcloud

        slide_5 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_5.placeholders[0].text = 'Top corporations, organizations, and nonprofits mentioned in social media posts referencing the coronavirus topic (previous 48 hours)'
        img_path = 'entity_wc.png'
        e_cloud = slide_5.shapes.add_picture(img_path, left=Inches(1.32), top=Inches(1.81), height=Inches(5.68),
                                             width=Inches(10.7))
        ## P.9 Headline Wordcloud ##

        slide_9 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_9.placeholders[0].text = 'Top phrases in headlines of recent U.S. articles about the coronavirus topic (previous 48 hours)'
        img_path = 'bigram.png'
        bigram = slide_9.shapes.add_picture(img_path, left=Inches(1.32), top=Inches(1.42), height=Inches(5.68),
                                                    width=Inches(10.7))

        # P.6 Entities Table
        slide_6 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_6.placeholders[0].text = 'Top 10 entities (people, products, and companies) mentioned in social media posts referencing COVID-19 (previous 48 hours)'

        ## table
        ent_table = slide_6.shapes.add_table(11, 2, width=Inches(6.12), height=Inches(5), left=Inches(3.61),
                                             top=Inches(1.57))

        ent_table.table.cell(0, 0).text = 'Ranking'
        cell = ent_table.table.rows[0].cells[0]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        ent_table.table.cell(0, 1).text = "Entity"
        cell = ent_table.table.rows[0].cells[1]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]


        for i in range(1, 11):
            for j in range(0, 2):
                cell = ent_table.table.rows[i].cells[j]
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text = str(self.entity_table.iloc[i - 1, j])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph = cell.text_frame.paragraphs[0]


        ## P.7 Engagement articles ##

        slide_7 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_7.placeholders[0].text = 'Top 5 U.S. articles about the coronavirus topic by social media engagement (previous 48 hours)'

        ## table
        eng_table = slide_7.shapes.add_table(6, 5, width=Inches(10.03), height=Inches(5), left=Inches(1.65),
                                             top=Inches(1.55))

        eng_table.table.cell(0, 0).text = 'Headline'
        cell = eng_table.table.rows[0].cells[0]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        eng_table.table.cell(0, 1).text = "Domain"
        cell = eng_table.table.rows[0].cells[1]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        eng_table.table.cell(0, 2).text = "Author"
        cell = eng_table.table.rows[0].cells[2]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        eng_table.table.cell(0, 3).text = "Facebook Interactions"
        cell = eng_table.table.rows[0].cells[3]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        eng_table.table.cell(0, 4).text = "Link"
        cell = eng_table.table.rows[0].cells[4]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        self.article_df['Facebook Interaction'] = self.article_df['Facebook Interaction'].apply(lambda x: self.human_format(x))

        for i in range(1, 6):
            for j in range(0, 5):
                cell = eng_table.table.rows[i].cells[j]
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text = [str(i) + '. ' + str(self.article_df.iloc[i - 1, j]) if j == 0 else str(self.article_df.iloc[i - 1, j])][0]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph = cell.text_frame.paragraphs[0]


        ## P.8 Publisher ##

        slide_8 = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
        slide_8.placeholders[0].text = 'Top 10 domains for articles about the coronavirus topic by volume of Facebook interactions (previous 48 hours)'

        ## table
        pub_table = slide_8.shapes.add_table(11, 2, width=Inches(6.12), height=Inches(5.37), left=Inches(3.61),
                                             top=Inches(1.57))

        pub_table.table.cell(0, 0).text = 'Domain'
        cell = pub_table.table.rows[0].cells[0]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        pub_table.table.cell(0, 1).text = "Facebook Interactions"
        cell = pub_table.table.rows[0].cells[1]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        self._set_cell_border(cell)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph = cell.text_frame.paragraphs[0]

        self.publisher_df['Facebook Interactions'] = self.publisher_df['Facebook Interactions'].apply(lambda x: self.human_format(x))

        for i in range(1, 11):
            for j in range(0, 2):
                cell = pub_table.table.rows[i].cells[j]
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_cell_border(cell)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text = str(self.publisher_df.iloc[i - 1, j])
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph = cell.text_frame.paragraphs[0]

        ppt_name = 'COVID-19 Monitoring Report ' + str(date.today()) + '.pptx'
        prs.save(ppt_name)


if __name__ == '__main__':
    covid = COVID_PPT()
    covid.PPT()
